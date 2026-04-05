"""
Slide Deck Generator MCP Server
Offers Slide Creation via FastMCP framework.
Stateful handling of slide structures.
"""

import os
import logging
from pathlib import Path
from typing import List

from fastmcp import FastMCP
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

logging.basicConfig(level=logging.INFO, format="[SLIDES-MCP] %(asctime)s | %(levelname)s - %(message)s")
logger = logging.getLogger("slides_mcp")

mcp_app = FastMCP(name="slide-builder-mcp")

# Storage for the current PPT session
session_state = {
    "deck": None,
    "output_file": None,
    "total_slides": 0,
}

# Theme Definitions
THEMES = {
    "professional": {
        "bg": RGBColor(255, 255, 255),
        "dark": RGBColor(15, 23, 42),
        "accent": RGBColor(37, 99, 235),
        "text": RGBColor(51, 65, 85),
        "font": "Arial"
    },
    "normal": {
        "bg": RGBColor(245, 245, 245),
        "dark": RGBColor(60, 60, 60),
        "accent": RGBColor(100, 100, 100),
        "text": RGBColor(30, 30, 30),
        "font": "Calibri"
    },
    "aesthetic": {
        "bg": RGBColor(253, 242, 248), # Soft Pinkish White
        "dark": RGBColor(80, 7, 36),   # Deep Maroon/Pink
        "accent": RGBColor(219, 39, 119), # Vibrant Pink
        "text": RGBColor(131, 24, 67),  # Darker Pink
        "font": "Georgia" # Elegant Serif
    }
}

DIMENSION_W = Inches(13.33)
DIMENSION_H = Inches(7.5)


def apply_bg(slide_obj, color_val: RGBColor):
    """Set the slide background color"""
    bg_fill = slide_obj.background.fill
    bg_fill.solid()
    bg_fill.fore_color.rgb = color_val


def insert_textbox(slide_obj, x: float, y: float, w: float, h: float, content: str, size: int, is_bold: bool, clr: RGBColor, font_name="Arial", alignment=PP_ALIGN.LEFT):
    """Insert a text shape into the slide"""
    box = slide_obj.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    paragraph.alignment = alignment
    run = paragraph.add_run()
    run.text = content
    run.font.size = Pt(size)
    run.font.bold = is_bold
    run.font.color.rgb = clr
    run.font.name = font_name


def create_the_title_slide(deck: Presentation, heading: str, theme_name="professional"):
    """Generate the initial cover slide with theme-specific styling"""
    theme = THEMES.get(theme_name, THEMES["professional"])
    empty_layout = deck.slide_layouts[6]
    cover = deck.slides.add_slide(empty_layout)
    apply_bg(cover, theme["dark"])
    
    if theme_name == "professional":
        accent_bar = cover.shapes.add_shape(1, Inches(0.5), Inches(2.5), Inches(0.1), Inches(2.5))
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = theme["accent"]
        accent_bar.line.fill.background()
        insert_textbox(cover, 0.8, 2.5, 11.5, 2.0, heading, 54, True, RGBColor(255, 255, 255), theme["font"])
    elif theme_name == "aesthetic":
        # Centered aesthetic title
        insert_textbox(cover, 0, 3.0, 13.33, 2.0, heading, 60, True, RGBColor(255, 240, 245), theme["font"], PP_ALIGN.CENTER)
    else:
        # Simple normal title
        insert_textbox(cover, 1.0, 3.0, 11.0, 2.0, heading, 44, True, RGBColor(255, 255, 255), theme["font"])

    subtitle = "Strategic Overview" if theme_name == "professional" else "Presentation"
    insert_textbox(cover, 0.8 if theme_name != "aesthetic" else 0, 4.8, 13.33 if theme_name == "aesthetic" else 8.0, 1.0, subtitle, 20, False, theme["accent"], theme["font"], PP_ALIGN.CENTER if theme_name == "aesthetic" else PP_ALIGN.LEFT)


def create_a_bullet_slide(deck: Presentation, heading: str, bullet_list: List[str], index: int, theme_name="professional"):
    """Generate a themed content slide"""
    theme = THEMES.get(theme_name, THEMES["professional"])
    empty_layout = deck.slide_layouts[6]
    slide = deck.slides.add_slide(empty_layout)
    apply_bg(slide, theme["bg"])

    if theme_name == "professional":
        header_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), DIMENSION_W, Inches(0.08))
        header_bar.fill.solid()
        header_bar.fill.fore_color.rgb = theme["accent"]
        header_bar.line.fill.background()
        insert_textbox(slide, 0.5, 0.4, 12.0, 1.0, heading, 32, True, theme["dark"], theme["font"])
        line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.01))
        line.fill.solid()
        line.line.color.rgb = RGBColor(226, 232, 240)
    elif theme_name == "aesthetic":
        insert_textbox(slide, 0.5, 0.5, 12.0, 1.0, heading, 36, True, theme["dark"], theme["font"])
        # Bottom aesthetic accent
        accent = slide.shapes.add_shape(1, Inches(0), Inches(7.3), DIMENSION_W, Inches(0.2))
        accent.fill.solid()
        accent.fill.fore_color.rgb = theme["accent"]
        accent.line.fill.background()
    else:
        insert_textbox(slide, 0.5, 0.5, 12.0, 1.0, heading, 28, True, theme["dark"], theme["font"])

    start_y = 1.8
    for idx, item in enumerate(bullet_list[:5]):
        # Themed bullet
        bullet_shape = 9 if theme_name == "professional" else (1 if theme_name == "aesthetic" else 1)
        bx = slide.shapes.add_shape(bullet_shape, Inches(0.6), Inches(start_y + idx * 0.9 + 0.15), Inches(0.12), Inches(0.12))
        bx.fill.solid()
        bx.fill.fore_color.rgb = theme["accent"]
        bx.line.fill.background()
        insert_textbox(slide, 0.9, start_y + idx * 0.9, 11.8, 0.8, item, 22 if theme_name != "normal" else 20, False, theme["text"], theme["font"])

    insert_textbox(slide, 12.0, 6.9, 1.0, 0.5, str(index), 12, False, theme["text"], theme["font"], PP_ALIGN.RIGHT)


@mcp_app.tool()
def init_presentation(filename: str, theme: str = "professional") -> str:
    """Prepare a new PPTX deck with a specific theme: 'professional', 'normal', or 'aesthetic'"""
    global session_state
    ppt = Presentation()
    ppt.slide_width = DIMENSION_W
    ppt.slide_height = DIMENSION_H

    session_state["deck"] = ppt
    session_state["output_file"] = filename
    session_state["total_slides"] = 0
    session_state["theme"] = theme if theme in THEMES else "professional"

    return f"New {session_state['theme']} presentation initialized for {filename}"


@mcp_app.tool()
def push_slide(title: str, bullet_points: List[str]) -> str:
    """Push a single slide into the active deck using the selected theme"""
    deck = session_state.get("deck")
    theme = session_state.get("theme", "professional")
    if not deck:
        return "Failure: Initialize presentation first."

    while len(bullet_points) < 3:
        bullet_points.append("Placeholder content.")
    bullet_points = bullet_points[:5]

    session_state["total_slides"] += 1
    current = session_state["total_slides"]

    if current == 1:
        create_the_title_slide(deck, title, theme)
    else:
        create_a_bullet_slide(deck, title, bullet_points, current - 1, theme)

    return f"Appended {theme} slide '{title}' (Slide {current})"


@mcp_app.tool()
def finalize_presentation() -> str:
    """Commit the slide deck to disk"""
    deck = session_state.get("deck")
    fname = session_state.get("output_file")
    if not deck or not fname:
        return "Failure: Nothing to save."
    fp = Path(os.getcwd()) / fname
    deck.save(str(fp))
    return f"Saved PPTX to {fp}"


if __name__ == "__main__":
    logger.info("Initializing Slide Builder MCP...")
    mcp_app.run(transport="stdio")
